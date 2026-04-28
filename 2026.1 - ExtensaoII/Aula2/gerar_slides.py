from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK_BG    = RGBColor(0x1A, 0x0A, 0x2E)   # roxo escuro
ACCENT     = RGBColor(0x8B, 0x5C, 0xF6)   # roxo claro
GOLD       = RGBColor(0xF5, 0xC5, 0x18)   # dourado
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG    = RGBColor(0x0D, 0x0D, 0x1A)
CODE_FG    = RGBColor(0xA8, 0xFF, 0x78)

blank_layout = prs.slide_layouts[6]  # completamente em branco


def new_slide():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BG
    return slide


def add_textbox(slide, text, left, top, width, height,
                font_size=24, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullets(slide, items, left, top, width, height,
                font_size=20, color=WHITE, title_color=GOLD,
                line_spacing=1.2):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        c = title_color if item.startswith("▶") else color
        run.font.color.rgb = c
        run.font.bold = item.startswith("▶")


def add_code_block(slide, code, left, top, width, height, font_size=14):
    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.word_wrap = True
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = CODE_BG
    tf = shape.text_frame
    tf.word_wrap = True
    first = True
    for line in code.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line if line else " "
        run.font.size = Pt(font_size)
        run.font.name = "Courier New"
        run.font.color.rgb = CODE_FG


def accent_bar(slide):
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(6.9), Inches(13.33), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def slide_number(slide, n):
    add_textbox(slide, str(n), 12.5, 7.1, 0.7, 0.3,
                font_size=11, color=ACCENT, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────
# SLIDE 1 – TÍTULO
# ─────────────────────────────────────────────
s = new_slide()
accent_bar(s)
add_textbox(s, "Introdução à Programação", 1, 1.5, 11, 1.2,
            font_size=44, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_textbox(s, "com Portugol Studio", 1, 2.6, 11, 0.9,
            font_size=32, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(s, "Do zero ao seu primeiro programa", 1, 3.6, 11, 0.7,
            font_size=20, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, "3 horas · Nível iniciante", 1, 5.8, 11, 0.5,
            font_size=15, color=ACCENT, align=PP_ALIGN.CENTER)
slide_number(s, 1)

# ─────────────────────────────────────────────
# SLIDE 2 – O QUE É PROGRAMAÇÃO?
# ─────────────────────────────────────────────
s = new_slide()
accent_bar(s)
add_textbox(s, "O que é programação?", 0.5, 0.4, 12, 0.9,
            font_size=32, bold=True, color=GOLD)
add_bullets(s, [
    "▶  Programar = dar instruções para o computador",
    "",
    "🍰  Receita de bolo:",
    "      1. Misture os ingredientes",
    "      2. Se a massa estiver muito grossa, adicione leite",
    "      3. Repita: bata por 5 minutos",
    "      4. Asse por 40 minutos",
    "",
    "💻  Código é exatamente isso — uma receita que o computador segue!",
], 0.8, 1.5, 11.5, 4.8, font_size=19)
slide_number(s, 2)

# ─────────────────────────────────────────────
# SLIDE 3 – PORTUGOL STUDIO
# ─────────────────────────────────────────────
s = new_slide()
accent_bar(s)
add_textbox(s, "Portugol Studio", 0.5, 0.4, 12, 0.9,
            font_size=32, bold=True, color=GOLD)
add_bullets(s, [
    "▶  O que é?",
    "      Ferramenta para aprender programação em português",
    "      Sintaxe simples, próxima do idioma natural",
    "",
    "▶  Por que usar?",
    "      Sem barreira de idioma",
    "      Feedback visual imediato (destaque de erros, console)",
    "      Gratuito e multiplataforma",
    "",
    "▶  Download",
    "      github.com/UNIVALI-LITE/Portugol-Studio/releases",
    "      Baixe o arquivo .jar (Windows/Linux) ou .dmg (Mac)",
], 0.8, 1.5, 11.5, 5.2, font_size=19)
slide_number(s, 3)

# ─────────────────────────────────────────────
# SLIDE 4 – PRIMEIRO PROGRAMA
# ─────────────────────────────────────────────
s = new_slide()
accent_bar(s)
add_textbox(s, "Primeiro Programa — Olá, Mundo!", 0.5, 0.4, 12, 0.9,
            font_size=28, bold=True, color=GOLD)
add_bullets(s, [
    "Todo programador começa aqui.",
    "Objetivo: exibir uma mensagem na tela.",
], 0.8, 1.4, 11.5, 0.9, font_size=18)

code = """\
programa
{
    funcao inicio()
    {
        escreva("Olá, Mundo!\\n")
    }
}"""
add_code_block(s, code, 0.8, 2.4, 7, 3.4, font_size=18)

add_bullets(s, [
    "▶  escreva()   →  mostra texto na tela",
    "▶  \\n          →  quebra de linha",
    "▶  funcao inicio()  →  ponto de entrada do programa",
], 8.1, 2.6, 4.8, 3, font_size=16)
slide_number(s, 4)

# ─────────────────────────────────────────────
# SLIDE 5 – VARIÁVEIS (CONCEITO)
# ─────────────────────────────────────────────
s = new_slide()
accent_bar(s)
add_textbox(s, "Variáveis", 0.5, 0.4, 12, 0.9,
            font_size=32, bold=True, color=GOLD)
add_bullets(s, [
    "▶  O que é uma variável?",
    "      Uma caixinha com nome que guarda um valor",
    "",
    "▶  Tipos de variáveis no Portugol:",
], 0.8, 1.4, 6, 2.8, font_size=19)

tabela = """\
  Tipo       Exemplo           Para que serve
 ──────────────────────────────────────────────
  inteiro    idade = 20        Números inteiros
  real       altura = 1.75     Números decimais
  cadeia     nome = "Ana"      Textos
  logico     ativo = verdadeiro  Verdadeiro/Falso"""
add_code_block(s, tabela, 0.8, 3.9, 11.5, 2.6, font_size=14)
slide_number(s, 5)

# ─────────────────────────────────────────────
# SLIDE 6 – VARIÁVEIS (CÓDIGO)
# ─────────────────────────────────────────────
s = new_slide()
accent_bar(s)
add_textbox(s, "Variáveis — na prática", 0.5, 0.4, 12, 0.9,
            font_size=28, bold=True, color=GOLD)

code = """\
programa
{
    funcao inicio()
    {
        inteiro idade = 20
        real    altura = 1.75
        cadeia  nome = "Ana"
        logico  estudante = verdadeiro

        escreva("Nome: ", nome, "\\n")
        escreva("Idade: ", idade, " anos\\n")
        escreva("Altura: ", altura, " m\\n")
        escreva("Estudante: ", estudante, "\\n")

        // Lendo do teclado
        escreva("Digite seu nome: ")
        leia(nome)
        escreva("Olá, ", nome, "!\\n")
    }
}"""
add_code_block(s, code, 0.8, 1.4, 7.5, 5.7, font_size=15)
add_bullets(s, [
    "▶  Declarar = criar a caixinha",
    "▶  Atribuir = colocar valor",
    "▶  leia()   = pegar valor do teclado",
    "",
    "💡  Nomes de variáveis:",
    "  • sem espaço",
    "  • sem acentos",
    "  • descritivos!",
], 8.7, 1.6, 4.2, 5, font_size=17)
slide_number(s, 6)

# ─────────────────────────────────────────────
# SLIDE 7 – CONDIÇÕES (CONCEITO)
# ─────────────────────────────────────────────
s = new_slide()
accent_bar(s)
add_textbox(s, "Condições — SE / SENÃO", 0.5, 0.4, 12, 0.9,
            font_size=32, bold=True, color=GOLD)
add_bullets(s, [
    "▶  O computador pode tomar decisões!",
    "",
    "🚦  Analogia do semáforo:",
    "      SE  a luz está verde",
    "          ➜  avance",
    "      SENÃO",
    "          ➜  pare",
    "",
    "▶  Estrutura no Portugol:",
], 0.8, 1.4, 6.2, 5.2, font_size=19)

code = """\
se (condição)
{
    // executa se VERDADEIRO
}
senao
{
    // executa se FALSO
}"""
add_code_block(s, code, 7.2, 1.6, 5.7, 3.2, font_size=17)
add_bullets(s, [
    "Operadores de comparação:",
    "  ==  igual a",
    "  !=  diferente de",
    "  >   maior que",
    "  <   menor que",
    "  >=  maior ou igual",
    "  <=  menor ou igual",
], 7.2, 5.0, 5.7, 2.2, font_size=15)
slide_number(s, 7)

# ─────────────────────────────────────────────
# SLIDE 8 – CONDIÇÕES (CÓDIGO)
# ─────────────────────────────────────────────
s = new_slide()
accent_bar(s)
add_textbox(s, "Condições — na prática", 0.5, 0.4, 12, 0.9,
            font_size=28, bold=True, color=GOLD)

code = """\
programa
{
    funcao inicio()
    {
        inteiro nota

        escreva("Digite sua nota (0-10): ")
        leia(nota)

        se (nota >= 6)
        {
            escreva("✓ Aprovado!\\n")
        }
        senao se (nota >= 4)
        {
            escreva("⚠ Recuperação\\n")
        }
        senao
        {
            escreva("✗ Reprovado\\n")
        }
    }
}"""
add_code_block(s, code, 0.8, 1.4, 7.5, 5.7, font_size=15)
add_bullets(s, [
    "▶  senao se",
    "   permite múltiplas",
    "   condições encadeadas",
    "",
    "💡  Fluxo de execução:",
    "   1. Testa nota >= 6",
    "   2. Se falso, testa >= 4",
    "   3. Se ainda falso,",
    "      executa o último senao",
    "",
    "   Apenas UM bloco",
    "   é executado!",
], 8.7, 1.6, 4.2, 5, font_size=17)
slide_number(s, 8)

# ─────────────────────────────────────────────
# SLIDE 9 – LOOPS (CONCEITO)
# ─────────────────────────────────────────────
s = new_slide()
accent_bar(s)
add_textbox(s, "Repetição — Loops", 0.5, 0.4, 12, 0.9,
            font_size=32, bold=True, color=GOLD)
add_bullets(s, [
    "▶  Repetição = executar o mesmo bloco várias vezes",
    "",
    "🏋️  Analogia da academia:",
    "      ENQUANTO  você tiver energia",
    "          faça 1 repetição",
    "",
    "📋  Ou com número fixo:",
    "      PARA  cada série de 1 até 3",
    "          faça 10 repetições",
], 0.8, 1.4, 6.5, 4.5, font_size=19)

code = """\
// Loop com condição
enquanto (condição)
{
    // repete enquanto for verdadeiro
}

// Loop com contador
para (inteiro i = 1; i <= 10; i++)
{
    // repete 10 vezes
}"""
add_code_block(s, code, 7.0, 1.6, 5.9, 4, font_size=16)
slide_number(s, 9)

# ─────────────────────────────────────────────
# SLIDE 10 – LOOPS (CÓDIGO)
# ─────────────────────────────────────────────
s = new_slide()
accent_bar(s)
add_textbox(s, "Loops — na prática", 0.5, 0.4, 12, 0.9,
            font_size=28, bold=True, color=GOLD)

code = """\
programa
{
    funcao inicio()
    {
        // Contagem regressiva com enquanto
        inteiro contador = 10
        enquanto (contador > 0)
        {
            escreva(contador, "...\\n")
            contador = contador - 1
        }
        escreva("Lançar!\\n")

        // Tabuada do 5 com para
        para (inteiro i = 1; i <= 10; i++)
        {
            escreva("5 x ", i, " = ", 5*i, "\\n")
        }
    }
}"""
add_code_block(s, code, 0.8, 1.4, 7.5, 5.7, font_size=14)
add_bullets(s, [
    "▶  enquanto",
    "   ideal quando não sei",
    "   quantas vezes repetir",
    "",
    "▶  para",
    "   ideal quando tenho",
    "   número fixo de vezes",
    "",
    "⚠  Cuidado com",
    "   loop infinito!",
    "   A condição PRECISA",
    "   se tornar falsa.",
], 8.7, 1.6, 4.2, 5, font_size=17)
slide_number(s, 10)

# ─────────────────────────────────────────────
# SLIDE 11 – LÓGICA BOOLEANA
# ─────────────────────────────────────────────
s = new_slide()
accent_bar(s)
add_textbox(s, "Lógica Booleana", 0.5, 0.4, 12, 0.9,
            font_size=32, bold=True, color=GOLD)
add_bullets(s, [
    "▶  Só existem dois valores: verdadeiro  e  falso",
    "",
    "💡  Usamos lógica booleana o tempo todo:",
    "",
    '  "Está chovendo  E  eu tenho guarda-chuva?"',
    "       ➜ verdadeiro E verdadeiro  =  VERDADEIRO  → sair",
    "",
    '  "Tenho dinheiro  OU  tenho crédito?"',
    "       ➜ falso OU verdadeiro      =  VERDADEIRO  → comprar",
    "",
    '  "NÃO estou com fome"',
    "       ➜ nao verdadeiro           =  FALSO",
], 0.8, 1.4, 11.5, 5.5, font_size=18)
slide_number(s, 11)

# ─────────────────────────────────────────────
# SLIDE 12 – AND / OR NO CÓDIGO
# ─────────────────────────────────────────────
s = new_slide()
accent_bar(s)
add_textbox(s, "Lógica Booleana — no Portugol", 0.5, 0.4, 12, 0.9,
            font_size=28, bold=True, color=GOLD)

code = """\
programa
{
    funcao inicio()
    {
        inteiro idade
        logico temCarteira

        escreva("Sua idade: ")
        leia(idade)
        escreva("Tem carteira de motorista? (verdadeiro/falso): ")
        leia(temCarteira)

        // E (AND) — ambas precisam ser verdadeiras
        se (idade >= 18 e temCarteira)
        {
            escreva("Pode dirigir!\\n")
        }

        // OU (OR) — pelo menos uma verdadeira
        se (idade < 18 ou nao temCarteira)
        {
            escreva("Não pode dirigir ainda.\\n")
        }
    }
}"""
add_code_block(s, code, 0.8, 1.4, 8.2, 5.7, font_size=13)

tabela = """\
 A      B      A e B   A ou B
 ────────────────────────────
 V      V        V       V
 V      F        F       V
 F      V        F       V
 F      F        F       F"""
add_code_block(s, tabela, 9.2, 2.0, 3.8, 2.6, font_size=13)
add_bullets(s, [
    "e   → AND (as duas)",
    "ou  → OR  (pelo menos 1)",
    "nao → NOT (inverte)",
], 9.2, 4.8, 3.8, 1.5, font_size=16)
slide_number(s, 12)

# ─────────────────────────────────────────────
# SLIDE 13 – COMBINANDO TUDO
# ─────────────────────────────────────────────
s = new_slide()
accent_bar(s)
add_textbox(s, "Combinando tudo — Exemplo Final", 0.5, 0.4, 12, 0.9,
            font_size=28, bold=True, color=GOLD)

code = """\
programa
{
    funcao inicio()
    {
        cadeia  nome
        inteiro tentativas = 0
        inteiro acertos    = 0
        cadeia  resposta

        escreva("Bem-vindo ao Quiz!\\n")
        escreva("Digite seu nome: ")
        leia(nome)

        enquanto (tentativas < 3)
        {
            escreva("\\nPergunta ", tentativas+1,
                    ": Capital do Brasil? ")
            leia(resposta)
            tentativas = tentativas + 1

            se (resposta == "Brasilia" ou resposta == "brasilia")
            {
                escreva("Correto!\\n")
                acertos = acertos + 1
            }
            senao
            {
                escreva("Errado! A resposta é Brasilia.\\n")
            }
        }
        escreva("\\n", nome, " acertou ", acertos,
                " de 3 perguntas!\\n")
    }
}"""
add_code_block(s, code, 0.5, 1.4, 12.3, 5.7, font_size=12)
slide_number(s, 13)

# ─────────────────────────────────────────────
# SLIDE 14 – PROJETO
# ─────────────────────────────────────────────
s = new_slide()
accent_bar(s)
add_textbox(s, "Projeto — Jogo de Adivinhação 🎮", 0.5, 0.4, 12, 0.9,
            font_size=28, bold=True, color=GOLD)
add_bullets(s, [
    "▶  Objetivo",
    "      Criar um jogo onde o computador 'pensa' em um número",
    "      de 1 a 100 e o jogador tenta adivinhar.",
    "",
    "▶  Regras do jogo",
    "      • A cada chute, o programa diz: 'maior' ou 'menor'",
    "      • Conta quantas tentativas o jogador usou",
    "      • Parabeniza quando acertar",
    "      • Permite jogar de novo",
    "",
    "▶  Conceitos que vão usar",
    "      variáveis  •  condições SE/SENÃO  •  loop enquanto",
    "      lógica booleana  •  entrada do teclado",
    "",
    "▶  Próxima aula: apresentação dos jogos + dicas de melhoria",
], 0.8, 1.4, 11.5, 5.5, font_size=18)
slide_number(s, 14)

# ─────────────────────────────────────────────
# SLIDE 15 – ENCERRAMENTO
# ─────────────────────────────────────────────
s = new_slide()
accent_bar(s)
add_textbox(s, "Resumo de hoje", 0.5, 0.4, 12, 0.9,
            font_size=32, bold=True, color=GOLD)
add_bullets(s, [
    "✅  Portugol Studio — ferramenta instalada e funcionando",
    "✅  Variáveis       — guardar e manipular dados",
    "✅  SE / SENÃO      — tomar decisões",
    "✅  Loops           — repetir ações",
    "✅  Lógica Booleana — verdadeiro, falso, E, OU, NÃO",
    "",
    "🚀  Para a próxima aula:",
    "      Desenvolver o Jogo de Adivinhação",
    "      Pode usar tudo que aprendemos hoje!",
    "",
    "      Dúvidas? Experimenta, erra, conserta — é assim que se aprende!",
], 0.8, 1.5, 11.5, 5, font_size=20)
slide_number(s, 15)

# ─────────────────────────────────────────────
# SALVAR
# ─────────────────────────────────────────────
out = "slides_intro_programacao.pptx"
prs.save(out)
print(f"Slides gerados: {out}  ({prs.slides.__len__()} slides)")
